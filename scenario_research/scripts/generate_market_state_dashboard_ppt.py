from __future__ import annotations

import argparse
import csv
import json
import re
from collections import Counter, defaultdict
from pathlib import Path

try:
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from PIL import Image, ImageDraw, ImageFilter
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
except ModuleNotFoundError as exc:
    missing_name = exc.name or "presentation dependency"
    raise SystemExit(
        "Missing presentation dependency "
        f"`{missing_name}`. Install `matplotlib`, `pillow`, and `python-pptx` to generate the deck."
    ) from exc


ROOT = Path(__file__).resolve().parents[1]
PLAN_PATH = ROOT / "plans" / "05_market_state_research_presentation_clean_20260417.md"
OUTPUT_DIR = ROOT / "outputs" / "presentations"
ASSET_DIR = OUTPUT_DIR / "_build" / "market_state_dashboard_20260417"
PPT_PATH = OUTPUT_DIR / "HedgeMate_Market_State_Research_Dashboard_20260417.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.42
GUTTER = 0.14
COL_W = (SLIDE_W - (MARGIN * 2) - (GUTTER * 11)) / 12

FONT_PRIMARY = "Pretendard"
FONT_LATIN = "Inter"

BG = "0B0F1A"
CARD = "132033"
CARD_ALT = "101A2C"
CARD_STROKE = "2B3854"
TEXT = "F5F8FF"
TEXT_SUB = "A0A7B8"
GREEN = "00FFB2"
MAGENTA = "FF2E9F"
RED = "FF3B3B"
CYAN = "4EE1FF"
AMBER = "F6B76F"
SOFT = "1C263A"

SHORT_NAME = {
    "Soft Landing / Goldilocks": "Soft Landing",
    "Slowdown / Recession / Deflation Risk": "Slowdown",
    "Stagflation / Reinflation / Energy Shock": "Stagflation",
    "USD Strength / KRW Weakness": "KRW Weakness",
    "China / Trade Fragmentation Shock": "China Shock",
    "Higher-for-Longer / Long-Rate Shock": "Long-Rate Shock",
    "Acute Global Stress / Liquidity Crunch": "Global Stress",
}


def sanitize_text(text: str) -> str:
    text = text.replace("`", "")
    text = re.sub(r"\[(.*?)\]\(.*?\)", r"\1", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def rgb(hex_code: str) -> RGBColor:
    return RGBColor.from_string(hex_code.replace("#", ""))


def gx(col_index: float) -> int:
    return Inches(MARGIN + col_index * (COL_W + GUTTER))


def gw(col_span: float) -> int:
    return Inches((COL_W * col_span) + (GUTTER * (col_span - 1)))


def gy(value: float) -> int:
    return Inches(value)


def newest_file(folder: Path, pattern: str) -> Path:
    matches = list(folder.glob(pattern))
    if not matches:
        raise FileNotFoundError(f"No files matched {folder / pattern}")
    return max(matches, key=lambda path: path.stat().st_mtime)


def require_file(path: Path) -> Path:
    if not path.exists():
        raise FileNotFoundError(f"Required input file not found: {path}")
    return path


def run_id_from_state_path(path: Path) -> str:
    prefix = "scenario_state_daily_"
    suffix = ".csv"
    name = path.name
    if not name.startswith(prefix) or not name.endswith(suffix):
        raise ValueError(f"Unexpected state filename format: {name}")
    return name[len(prefix) : -len(suffix)]


def format_anchor_short(metadata: dict[str, object]) -> str:
    anchor_date = str(metadata.get("anchor_date") or "")
    parts = anchor_date.split("-")
    if len(parts) == 3:
        return f"{parts[1]}.{parts[2]}"
    return anchor_date or "N/A"


def format_anchor_caption(metadata: dict[str, object]) -> str:
    anchor_date = str(metadata.get("anchor_date") or "")
    year = anchor_date.split("-")[0] if "-" in anchor_date else ""
    return f"{year} snapshot" if year else "snapshot"


def parse_markdown_sections(path: Path) -> dict[str, dict[str, list[str] | str]]:
    sections: dict[str, dict[str, list[str] | str]] = {}
    current_key: str | None = None

    for raw_line in path.read_text(encoding="utf-8").splitlines():
        line = raw_line.rstrip()
        stripped = line.strip()

        slide_match = re.match(r"## Slide\s+(\d+)\.\s+(.+)", stripped)
        if slide_match:
            current_key = f"slide_{slide_match.group(1)}"
            sections[current_key] = {"title": slide_match.group(2), "lines": []}
            continue

        if stripped.startswith("## Appendix."):
            current_key = "appendix"
            sections[current_key] = {"title": stripped.replace("## ", "", 1), "lines": []}
            continue

        if stripped.startswith("## References"):
            current_key = "references"
            sections[current_key] = {"title": "References", "lines": []}
            continue

        if stripped == "---":
            continue

        if current_key:
            sections[current_key]["lines"].append(line)

    return sections


def extract_bullets(lines: list[str]) -> list[str]:
    bullets: list[str] = []
    current: str | None = None

    for line in lines:
        stripped = line.strip()
        if stripped.startswith("- "):
            if current:
                bullets.append(sanitize_text(current))
            current = stripped[2:].strip()
            continue
        if stripped.endswith(":") or stripped.startswith("###") or re.match(r"\d+\.", stripped):
            if current:
                bullets.append(sanitize_text(current))
                current = None
            continue
        if current and stripped:
            current = f"{current} {stripped}"

    if current:
        bullets.append(sanitize_text(current))

    return bullets


def extract_phases(lines: list[str]) -> list[dict[str, list[str] | str]]:
    phases: list[dict[str, list[str] | str]] = []
    current: dict[str, list[str] | str] | None = None

    for line in lines:
        stripped = line.strip()
        if stripped.startswith("### Phase "):
            if current:
                phases.append(current)
            current = {"title": sanitize_text(stripped.replace("### ", "", 1)), "bullets": []}
            continue
        if stripped.startswith("- ") and current:
            current["bullets"].append(sanitize_text(stripped[2:].strip()))

    if current:
        phases.append(current)

    return phases


def load_data(run_id: str | None = None) -> dict[str, object]:
    processed_dir = ROOT / "outputs" / "processed"
    reports_dir = ROOT / "outputs" / "reports"

    if run_id:
        state_path = require_file(processed_dir / f"scenario_state_daily_{run_id}.csv")
    else:
        state_path = newest_file(processed_dir, "scenario_state_daily_*.csv")
        run_id = run_id_from_state_path(state_path)

    driver_path = require_file(reports_dir / f"scenario_driver_table_{run_id}.csv")
    metadata_path = require_file(reports_dir / f"scenario_snapshot_metadata_{run_id}.json")

    with state_path.open(encoding="utf-8", newline="") as file:
        state_rows = list(csv.DictReader(file))

    latest_date = max(row["date"] for row in state_rows)
    latest_rows = [row for row in state_rows if row["date"] == latest_date]
    for row in latest_rows:
        row["structured_score"] = float(row["structured_score"])
        row["confidence"] = float(row["confidence"])
        row["coverage_ratio"] = float(row["coverage_ratio"])
    latest_rows.sort(key=lambda row: row["structured_score"], reverse=True)

    with driver_path.open(encoding="utf-8", newline="") as file:
        driver_rows = list(csv.DictReader(file))

    latest_driver_rows = [row for row in driver_rows if row["date"] == latest_date]
    drivers_by_scenario: dict[str, list[dict[str, str]]] = defaultdict(list)
    for row in latest_driver_rows:
        drivers_by_scenario[row["scenario_name"]].append(row)
    for key in drivers_by_scenario:
        drivers_by_scenario[key].sort(key=lambda row: int(row["driver_rank"]))

    with metadata_path.open(encoding="utf-8") as file:
        metadata = json.load(file)

    return {
        "latest_date": latest_date,
        "latest_rows": latest_rows,
        "drivers_by_scenario": drivers_by_scenario,
        "metadata": metadata,
        "state_rows": state_rows,
    }


def short_name(name: str) -> str:
    return SHORT_NAME.get(name, name)


def state_color(state_label: str, scenario_name: str) -> str:
    if state_label == "ACTIVE":
        return GREEN
    if state_label == "WATCH":
        return MAGENTA
    if "Stress" in scenario_name or "Shock" in scenario_name:
        return RED
    return "54627E"


def ensure_dirs() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def add_glow(base: Image.Image, center: tuple[int, int], radius: int, color_hex: str, alpha: int) -> None:
    overlay = Image.new("RGBA", base.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)
    color = tuple(int(color_hex[i : i + 2], 16) for i in (0, 2, 4))
    x, y = center
    draw.ellipse((x - radius, y - radius, x + radius, y + radius), fill=color + (alpha,))
    overlay = overlay.filter(ImageFilter.GaussianBlur(radius // 2))
    base.alpha_composite(overlay)


def create_background(path: Path, variant: str) -> None:
    width, height = 1600, 900
    img = Image.new("RGBA", (width, height), tuple(int(BG[i : i + 2], 16) for i in (0, 2, 4)) + (255,))
    draw = ImageDraw.Draw(img, "RGBA")

    for x in range(0, width, 86):
        draw.line((x, 0, x, height), fill=(54, 67, 97, 24), width=1)
    for y in range(0, height, 86):
        draw.line((0, y, width, y), fill=(54, 67, 97, 18), width=1)
    for offset in range(-height, width, 160):
        draw.line((offset, 0, offset + height, height), fill=(38, 49, 76, 14), width=1)

    if variant == "hero":
        add_glow(img, (1250, 180), 230, GREEN, 118)
        add_glow(img, (260, 720), 280, MAGENTA, 64)
    elif variant == "analysis":
        add_glow(img, (1380, 230), 220, CYAN, 86)
        add_glow(img, (250, 240), 240, MAGENTA, 60)
    elif variant == "roadmap":
        add_glow(img, (1320, 180), 250, MAGENTA, 88)
        add_glow(img, (250, 760), 220, GREEN, 72)
    else:
        add_glow(img, (1320, 160), 250, GREEN, 82)
        add_glow(img, (160, 740), 220, MAGENTA, 58)

    vignette = Image.new("RGBA", img.size, (0, 0, 0, 0))
    vdraw = ImageDraw.Draw(vignette)
    vdraw.rectangle((0, 0, width, height), fill=(0, 0, 0, 52))
    vignette = vignette.filter(ImageFilter.GaussianBlur(24))
    img.alpha_composite(vignette)
    img.save(path)


def create_bar_chart(latest_rows: list[dict[str, object]], path: Path) -> None:
    ordered = list(reversed(latest_rows))
    labels = [short_name(str(row["scenario_name"])) for row in ordered]
    values = [float(row["structured_score"]) for row in ordered]
    colors = [f"#{state_color(str(row['state_label']), str(row['scenario_name']))}" for row in ordered]

    fig, ax = plt.subplots(figsize=(7.0, 3.8), facecolor=(0, 0, 0, 0))
    ax.set_facecolor((0, 0, 0, 0))
    ax.barh(labels, values, color=colors, height=0.58)
    ax.set_xlim(0, 100)
    ax.grid(axis="x", color="#42506F", alpha=0.22, linewidth=0.8)
    ax.tick_params(axis="x", colors="#7C879C", labelsize=9)
    ax.tick_params(axis="y", colors="#EEF4FF", labelsize=11)

    for spine in ax.spines.values():
        spine.set_visible(False)

    for idx, value in enumerate(values):
        ax.text(min(value + 1.2, 98), idx, f"{value:.1f}", va="center", ha="left", color="#F5F8FF", fontsize=11, weight="bold")

    plt.tight_layout(pad=0.6)
    fig.savefig(path, dpi=220, transparent=True)
    plt.close(fig)


def create_donut_chart(latest_rows: list[dict[str, object]], path: Path) -> None:
    values = [float(row["structured_score"]) for row in latest_rows]
    palette = [f"#{GREEN}", f"#{MAGENTA}", f"#{CYAN}", "#8995AE", f"#{RED}", "#626E85", "#3D4860"]

    fig, ax = plt.subplots(figsize=(3.8, 3.8), facecolor=(0, 0, 0, 0))
    ax.set_facecolor((0, 0, 0, 0))
    ax.pie(
        values,
        startangle=90,
        colors=palette[: len(values)],
        wedgeprops={"width": 0.32, "edgecolor": f"#{BG}", "linewidth": 2},
    )
    ax.text(0, 0.10, "Regime", ha="center", va="center", color="#A0A7B8", fontsize=12)
    ax.text(0, -0.08, "Mix", ha="center", va="center", color="#F5F8FF", fontsize=23, weight="bold")
    ax.set(aspect="equal")
    plt.tight_layout(pad=0.3)
    fig.savefig(path, dpi=220, transparent=True)
    plt.close(fig)


def create_history_chart(state_rows: list[dict[str, str]], latest_rows: list[dict[str, object]], path: Path) -> None:
    top_names = [str(row["scenario_name"]) for row in latest_rows[:3]]
    series_by_name: dict[str, list[tuple[str, float]]] = defaultdict(list)
    for row in state_rows:
        if row["scenario_name"] in top_names:
            series_by_name[row["scenario_name"]].append((row["date"], float(row["structured_score"])))

    fig, ax = plt.subplots(figsize=(7.2, 2.4), facecolor=(0, 0, 0, 0))
    ax.set_facecolor((0, 0, 0, 0))

    colors = [f"#{GREEN}", f"#{MAGENTA}", f"#{CYAN}"]
    for color, name in zip(colors, top_names):
        points = sorted(series_by_name[name], key=lambda item: item[0])[-45:]
        values = [score for _, score in points]
        ax.plot(values, linewidth=2.2, color=color, label=short_name(name))
        ax.fill_between(range(len(values)), values, [0] * len(values), color=color, alpha=0.05)

    ax.set_ylim(0, 85)
    ax.set_xlim(0, 44)
    ax.grid(axis="y", color="#42506F", alpha=0.18, linewidth=0.8)
    ax.tick_params(axis="x", colors="#6E7890", labelsize=8)
    ax.tick_params(axis="y", colors="#6E7890", labelsize=8)
    ax.set_xticks([0, 11, 22, 33, 44])
    ax.set_xticklabels(["-45d", "-34d", "-23d", "-12d", "Latest"])
    for spine in ax.spines.values():
        spine.set_visible(False)

    legend = ax.legend(loc="upper left", frameon=False, fontsize=9)
    for text in legend.get_texts():
        text.set_color("#F5F8FF")

    plt.tight_layout(pad=0.3)
    fig.savefig(path, dpi=220, transparent=True)
    plt.close(fig)


def add_background(slide, image_path: Path) -> None:
    slide.shapes.add_picture(str(image_path), 0, 0, width=Inches(SLIDE_W), height=Inches(SLIDE_H))


def style_run(font, size: int, color_hex: str, *, bold: bool = False, font_name: str = FONT_PRIMARY) -> None:
    font.name = font_name
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = rgb(color_hex)


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text: str,
    *,
    size: int,
    color_hex: str = TEXT,
    bold: bool = False,
    font_name: str = FONT_PRIMARY,
    align=PP_ALIGN.LEFT,
    line_spacing: float = 1.1,
) -> None:
    textbox = slide.shapes.add_textbox(x, y, w, h)
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.text = text
    for paragraph in frame.paragraphs:
        paragraph.alignment = align
        paragraph.line_spacing = line_spacing
        for run in paragraph.runs:
            style_run(run.font, size, color_hex, bold=bold, font_name=font_name)


def add_paragraph_block(
    slide,
    x,
    y,
    w,
    h,
    items: list[str],
    *,
    size: int,
    color_hex: str = TEXT_SUB,
    bullet: bool = False,
) -> None:
    textbox = slide.shapes.add_textbox(x, y, w, h)
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0

    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.line_spacing = 1.18
        paragraph.space_after = Pt(8)
        paragraph.text = f"• {item}" if bullet else item
        style_run(paragraph.runs[0].font, size, color_hex, font_name=FONT_PRIMARY)


def add_glass_card(slide, x, y, w, h, *, fill_hex: str = CARD, transparency: float = 0.18, line_hex: str = CARD_STROKE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb(fill_hex)
    fill.transparency = transparency
    shape.line.color.rgb = rgb(line_hex)
    shape.line.width = Pt(1)
    return shape


def add_kpi_card(slide, x, y, w, h, label: str, value: str, caption: str, accent_hex: str) -> None:
    add_glass_card(slide, x, y, w, h)
    add_textbox(slide, x + Inches(0.18), y + Inches(0.16), w - Inches(0.32), Inches(0.22), label.upper(), size=9, color_hex=TEXT_SUB, bold=True, font_name=FONT_LATIN)
    value_size = 18 if len(value) >= 7 else 22
    add_textbox(slide, x + Inches(0.18), y + Inches(0.42), w - Inches(0.32), Inches(0.42), value, size=value_size, color_hex=TEXT, bold=True, font_name=FONT_LATIN)
    add_textbox(slide, x + Inches(0.18), y + Inches(0.84), w - Inches(0.32), Inches(0.24), caption, size=10, color_hex=accent_hex)


def add_chip(slide, x, y, text: str, accent_hex: str, *, width_in: float = 1.35, active: bool = True) -> None:
    fill_hex = accent_hex if active else SOFT
    border_hex = accent_hex if active else CARD_STROKE
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, Inches(width_in), Inches(0.34))
    chip.fill.solid()
    chip.fill.fore_color.rgb = rgb(fill_hex)
    chip.fill.transparency = 0.80 if active else 0.45
    chip.line.color.rgb = rgb(border_hex)
    chip.line.width = Pt(1)
    add_textbox(slide, x + Inches(0.08), y + Inches(0.05), Inches(width_in - 0.16), Inches(0.22), text, size=9, color_hex=BG if active else TEXT_SUB, bold=True, font_name=FONT_LATIN, align=PP_ALIGN.CENTER)


def add_gauge(slide, x, y, w, label: str, value: float, state_label: str, accent_hex: str) -> None:
    add_textbox(slide, x, y, w, Inches(0.18), label, size=11, color_hex=TEXT_SUB)
    add_textbox(slide, x + w - Inches(0.90), y, Inches(0.90), Inches(0.18), f"{value:.1f}", size=11, color_hex=TEXT, bold=True, align=PP_ALIGN.RIGHT, font_name=FONT_LATIN)

    rail = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y + Inches(0.24), w, Inches(0.14))
    rail.fill.solid()
    rail.fill.fore_color.rgb = rgb("273347")
    rail.line.fill.background()

    bar_width = max(0.35, min((value / 100) * 4.85, (w / 914400) - 0.08))
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y + Inches(0.24), Inches(bar_width), Inches(0.14))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(accent_hex)
    bar.fill.transparency = 0.18
    bar.line.fill.background()

    add_chip(slide, x + w - Inches(1.02), y + Inches(0.43), state_label, accent_hex, width_in=0.92, active=True)


def add_footer(slide, page_number: int) -> None:
    add_textbox(slide, gx(0), gy(7.04), gw(6), Inches(0.16), "HedgeMate | AI-powered portfolio hedging platform", size=9, color_hex="67738B", font_name=FONT_LATIN)
    add_textbox(slide, gx(10.4), gy(7.04), gw(1.6), Inches(0.16), f"{page_number:02d}", size=9, color_hex="67738B", font_name=FONT_LATIN, align=PP_ALIGN.RIGHT)


def slide_cover(prs: Presentation, sections, data, assets) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, assets["hero_bg"])
    latest_rows = data["latest_rows"]
    metadata = data["metadata"]
    state_counts = Counter(row["state_label"] for row in latest_rows)

    add_textbox(slide, gx(0), gy(0.40), gw(2.4), Inches(0.20), "HedgeMate", size=11, color_hex=GREEN, bold=True, font_name=FONT_LATIN)
    add_textbox(slide, gx(0), gy(0.72), gw(6.2), Inches(1.18), "시장 상태를 읽어주는\n설명 가능한 리서치 엔진", size=28, color_hex=TEXT, bold=True)
    add_textbox(slide, gx(0), gy(1.98), gw(5.2), Inches(0.54), "포트폴리오 분석과 리스크 관리 결과보다 먼저, 지금 시장이 어떤 환경인지와 그 해석 근거를 빠르게 보여주는 HedgeMate의 시장 상태 레이어입니다.", size=13, color_hex=TEXT_SUB)

    kpi_y = 0.52
    add_kpi_card(slide, gx(7.35), gy(kpi_y), Inches(1.35), Inches(1.18), "Anchor Date", format_anchor_short(metadata), format_anchor_caption(metadata), GREEN)
    add_kpi_card(slide, gx(8.88), gy(kpi_y), Inches(1.35), Inches(1.18), "Coverage", f"{metadata['anchor_ticker_count']} / {metadata['total_tickers']}", f"{metadata['anchor_ticker_coverage_ratio'] * 100:.1f}%", CYAN)
    add_kpi_card(slide, gx(10.41), gy(kpi_y), Inches(1.35), Inches(1.18), "Current Phase", "4 / 10", "summary layer", MAGENTA)
    add_kpi_card(slide, gx(11.94), gy(kpi_y), Inches(0.97), Inches(1.18), "States", str(len(latest_rows)), "tracked", TEXT_SUB)

    overview_bullets = [
        "HedgeMate는 포트폴리오 분석 및 리스크 관리 솔루션이다.",
        "시나리오 리서치는 시장 상태를 읽어주는 해석 레이어다.",
        "자산 분석 전에 현재 환경과 포트폴리오 영향의 맥락을 먼저 설명한다.",
    ]
    left_card = add_glass_card(slide, gx(0), gy(2.72), gw(5.6), Inches(3.62))
    add_textbox(slide, left_card.left + Inches(0.22), left_card.top + Inches(0.18), Inches(2.4), Inches(0.18), "WHY IT MATTERS", size=10, color_hex=TEXT_SUB, bold=True, font_name=FONT_LATIN)
    add_paragraph_block(slide, left_card.left + Inches(0.22), left_card.top + Inches(0.54), left_card.width - Inches(0.44), Inches(1.72), overview_bullets, size=12, color_hex=TEXT_SUB, bullet=True)
    add_textbox(slide, left_card.left + Inches(0.22), left_card.top + Inches(2.34), Inches(2.2), Inches(0.18), "CORE ROLES", size=10, color_hex=GREEN, bold=True, font_name=FONT_LATIN)
    add_paragraph_block(slide, left_card.left + Inches(0.22), left_card.top + Inches(2.62), left_card.width - Inches(0.44), Inches(0.78), extract_bullets(sections["slide_1"]["lines"])[3:], size=12, color_hex=TEXT, bullet=True)

    right_card = add_glass_card(slide, gx(5.9), gy(2.18), gw(6.1), Inches(4.16), fill_hex=CARD_ALT)
    add_textbox(slide, right_card.left + Inches(0.22), right_card.top + Inches(0.18), Inches(3.6), Inches(0.18), "LATEST REGIME STACK", size=10, color_hex=TEXT_SUB, bold=True, font_name=FONT_LATIN)
    add_chip(slide, right_card.left + Inches(0.22), right_card.top + Inches(0.48), "ACTIVE", GREEN, width_in=0.88)
    add_chip(slide, right_card.left + Inches(1.16), right_card.top + Inches(0.48), "WATCH", MAGENTA, width_in=0.88)
    add_chip(slide, right_card.left + Inches(2.10), right_card.top + Inches(0.48), "OFF", RED, width_in=0.68, active=False)

    gauge_top = right_card.top + Inches(0.96)
    for idx, row in enumerate(latest_rows[:5]):
        add_gauge(
            slide,
            right_card.left + Inches(0.22),
            gauge_top + Inches(idx * 0.58),
            right_card.width - Inches(0.44),
            short_name(str(row["scenario_name"])),
            float(row["structured_score"]),
            str(row["state_label"]),
            state_color(str(row["state_label"]), str(row["scenario_name"])),
        )

    add_textbox(slide, right_card.left + Inches(0.22), right_card.top + Inches(3.54), Inches(2.4), Inches(0.18), "CURRENT SNAPSHOT", size=10, color_hex=TEXT_SUB, bold=True, font_name=FONT_LATIN)
    snapshot = [f"ACTIVE {state_counts.get('ACTIVE', 0)}", f"WATCH {state_counts.get('WATCH', 0)}", f"OFF {state_counts.get('OFF', 0)}"]
    add_paragraph_block(slide, right_card.left + Inches(0.22), right_card.top + Inches(3.80), Inches(2.1), Inches(0.32), snapshot, size=11, color_hex=TEXT)
    add_footer(slide, 1)


def slide_need(prs: Presentation, sections, assets) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, assets["analysis_bg"])
    bullets = extract_bullets(sections["slide_2"]["lines"])

    add_textbox(slide, gx(0), gy(0.48), gw(5.8), Inches(0.30), "왜 시장 상태 레이어가 필요한가", size=24, color_hex=TEXT, bold=True)
    add_textbox(slide, gx(0), gy(0.88), gw(4.4), Inches(0.30), "Same asset, different regime, different risk meaning.", size=12, color_hex=GREEN, bold=True, font_name=FONT_LATIN)

    main_card = add_glass_card(slide, gx(0), gy(1.38), gw(5.25), Inches(4.92))
    add_textbox(slide, main_card.left + Inches(0.22), main_card.top + Inches(0.18), Inches(2.2), Inches(0.18), "KEY MESSAGE", size=10, color_hex=TEXT_SUB, bold=True, font_name=FONT_LATIN)
    add_paragraph_block(slide, main_card.left + Inches(0.22), main_card.top + Inches(0.50), main_card.width - Inches(0.44), Inches(1.70), bullets[:3], size=13, color_hex=TEXT_SUB, bullet=True)
    add_textbox(slide, main_card.left + Inches(0.22), main_card.top + Inches(2.48), Inches(2.3), Inches(0.18), "CORE QUESTIONS", size=10, color_hex=GREEN, bold=True, font_name=FONT_LATIN)
    add_paragraph_block(slide, main_card.left + Inches(0.22), main_card.top + Inches(2.78), main_card.width - Inches(0.44), Inches(1.20), bullets[3:], size=13, color_hex=TEXT, bullet=True)

    scenario_examples = [
        ("KRW Weakness", "같은 자산도 환율과 수입물가의 영향이 커집니다.", GREEN),
        ("Risk-Off", "상관관계가 바뀌고 헤지 자산의 역할이 더 중요해집니다.", MAGENTA),
        ("Rate Shock", "듀레이션과 성장주 멀티플 압박을 다르게 읽어야 합니다.", RED),
    ]
    card_x = gx(5.7)
    for idx, (label, body, accent) in enumerate(scenario_examples):
        y = gy(1.46 + idx * 1.56)
        card = add_glass_card(slide, card_x, y, gw(5.7), Inches(1.18), fill_hex=CARD_ALT)
        add_chip(slide, card.left + Inches(0.18), card.top + Inches(0.16), label, accent, width_in=1.5, active=True)
        add_textbox(slide, card.left + Inches(1.86), card.top + Inches(0.20), card.width - Inches(2.04), Inches(0.66), body, size=13, color_hex=TEXT)

    bottom = add_glass_card(slide, gx(5.7), gy(6.06), gw(5.7), Inches(0.76))
    add_textbox(slide, bottom.left + Inches(0.20), bottom.top + Inches(0.18), bottom.width - Inches(0.40), Inches(0.24), "해석이 빠진 리스크 수치는 설명력이 약하고, 시장 상태 레이어가 있어야 HedgeMate의 결과가 더 설득력 있게 보입니다.", size=12, color_hex=TEXT_SUB)
    add_footer(slide, 2)


def slide_engine(prs: Presentation, data, assets) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, assets["analysis_bg"])

    add_textbox(slide, gx(0), gy(0.48), gw(5.8), Inches(0.30), "현재 만들고 있는 방식", size=24, color_hex=TEXT, bold=True)
    add_textbox(slide, gx(0), gy(0.90), gw(5.2), Inches(0.24), "Market Data -> Scenario Score -> State Classification -> Explainable Summary", size=12, color_hex=GREEN, bold=True, font_name=FONT_LATIN)

    steps = [
        ("01", "시나리오 정의", "어떤 시장 상태를 추적할지 정리"),
        ("02", "정형 데이터 수집", "환율, 금리, 주가지수, 원자재, 변동성 수집"),
        ("03", "시나리오 점수 계산", "신호를 합쳐 강도를 계산"),
        ("04", "시장 상태 요약", "활성 시나리오와 근거를 설명"),
    ]

    step_width = Inches(2.88)
    for idx, (num, title, body) in enumerate(steps):
        x = gx(idx * 3.02)
        card = add_glass_card(slide, x, gy(1.56), step_width, Inches(1.52), fill_hex=CARD_ALT if idx % 2 else CARD)
        add_textbox(slide, card.left + Inches(0.16), card.top + Inches(0.16), Inches(0.42), Inches(0.20), num, size=12, color_hex=GREEN if idx in (0, 3) else MAGENTA, bold=True, font_name=FONT_LATIN)
        add_textbox(slide, card.left + Inches(0.16), card.top + Inches(0.48), card.width - Inches(0.32), Inches(0.28), title, size=16, color_hex=TEXT, bold=True)
        add_textbox(slide, card.left + Inches(0.16), card.top + Inches(0.86), card.width - Inches(0.32), Inches(0.46), body, size=11, color_hex=TEXT_SUB)

    trend_card = add_glass_card(slide, gx(0), gy(3.42), gw(7.0), Inches(2.52))
    add_textbox(slide, trend_card.left + Inches(0.20), trend_card.top + Inches(0.18), Inches(3.4), Inches(0.18), "RECENT REGIME INTENSITY TREND", size=10, color_hex=TEXT_SUB, bold=True, font_name=FONT_LATIN)
    slide.shapes.add_picture(str(assets["history_chart"]), trend_card.left + Inches(0.16), trend_card.top + Inches(0.44), width=trend_card.width - Inches(0.32), height=trend_card.height - Inches(0.60))

    evidence_card = add_glass_card(slide, gx(7.2), gy(3.42), gw(4.8), Inches(2.52), fill_hex=CARD_ALT)
    top_row = data["latest_rows"][0]
    drivers = data["drivers_by_scenario"][top_row["scenario_name"]][:3]
    add_textbox(slide, evidence_card.left + Inches(0.20), evidence_card.top + Inches(0.18), Inches(3.4), Inches(0.18), "LATEST ACTIVE REGIME", size=10, color_hex=GREEN, bold=True, font_name=FONT_LATIN)
    add_textbox(slide, evidence_card.left + Inches(0.20), evidence_card.top + Inches(0.46), Inches(3.6), Inches(0.30), short_name(str(top_row["scenario_name"])), size=18, color_hex=TEXT, bold=True)
    add_textbox(slide, evidence_card.left + Inches(0.20), evidence_card.top + Inches(0.82), Inches(3.6), Inches(0.20), f"Score {top_row['structured_score']:.2f} | Confidence {top_row['confidence']:.2f}", size=11, color_hex=TEXT_SUB, font_name=FONT_LATIN)

    driver_lines = [f"{item['signal_label']} | contribution {float(item['contribution']):+.3f}" for item in drivers]
    add_paragraph_block(slide, evidence_card.left + Inches(0.20), evidence_card.top + Inches(1.22), evidence_card.width - Inches(0.40), Inches(1.00), driver_lines, size=11, color_hex=TEXT, bullet=True)
    add_footer(slide, 3)


def slide_output(prs: Presentation, sections, data, assets) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, assets["default_bg"])

    add_textbox(slide, gx(0), gy(0.48), gw(6.2), Inches(0.30), "현재 출력 목표와 최신 스냅샷", size=24, color_hex=TEXT, bold=True)
    add_textbox(slide, gx(0), gy(0.90), gw(5.8), Inches(0.24), "Daily scenario state, scenario score, driver table, summary note.", size=12, color_hex=GREEN, bold=True, font_name=FONT_LATIN)

    chip_y = gy(1.30)
    for idx, row in enumerate(data["latest_rows"][:5]):
        add_chip(slide, gx(0) + Inches(idx * 1.42), chip_y, short_name(str(row["scenario_name"])), state_color(str(row["state_label"]), str(row["scenario_name"])), width_in=1.28, active=str(row["state_label"]) != "OFF")

    left_card = add_glass_card(slide, gx(0), gy(1.76), gw(7.0), Inches(4.68))
    add_textbox(slide, left_card.left + Inches(0.20), left_card.top + Inches(0.18), Inches(2.6), Inches(0.18), "SCENARIO SCOREBOARD", size=10, color_hex=TEXT_SUB, bold=True, font_name=FONT_LATIN)
    slide.shapes.add_picture(str(assets["bar_chart"]), left_card.left + Inches(0.18), left_card.top + Inches(0.38), width=left_card.width - Inches(0.36), height=left_card.height - Inches(0.56))

    donut_card = add_glass_card(slide, gx(7.2), gy(1.76), gw(4.8), Inches(2.22), fill_hex=CARD_ALT)
    add_textbox(slide, donut_card.left + Inches(0.18), donut_card.top + Inches(0.18), Inches(2.0), Inches(0.18), "REGIME MIX", size=10, color_hex=TEXT_SUB, bold=True, font_name=FONT_LATIN)
    slide.shapes.add_picture(str(assets["donut_chart"]), donut_card.left + Inches(0.10), donut_card.top + Inches(0.34), height=Inches(1.72))
    legend_y = donut_card.top + Inches(0.50)
    for idx, row in enumerate(data["latest_rows"][:4]):
        accent = state_color(str(row["state_label"]), str(row["scenario_name"]))
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, donut_card.left + Inches(2.18), legend_y + Inches(idx * 0.32), Inches(0.10), Inches(0.10))
        dot.fill.solid()
        dot.fill.fore_color.rgb = rgb(accent)
        dot.line.fill.background()
        add_textbox(slide, donut_card.left + Inches(2.34), legend_y + Inches(idx * 0.28) - Inches(0.02), Inches(2.16), Inches(0.18), short_name(str(row["scenario_name"])), size=10, color_hex=TEXT_SUB)

    driver_card = add_glass_card(slide, gx(7.2), gy(4.12), gw(4.8), Inches(2.48))
    add_textbox(slide, driver_card.left + Inches(0.18), driver_card.top + Inches(0.18), Inches(2.2), Inches(0.18), "TOP DRIVERS", size=10, color_hex=GREEN, bold=True, font_name=FONT_LATIN)
    top_drivers = data["drivers_by_scenario"][data["latest_rows"][0]["scenario_name"]][:3]
    driver_lines = [f"{item['signal_label']} | normalized {float(item['normalized_value']):+.3f}" for item in top_drivers]
    add_paragraph_block(slide, driver_card.left + Inches(0.18), driver_card.top + Inches(0.50), driver_card.width - Inches(0.36), Inches(1.08), driver_lines, size=11, color_hex=TEXT, bullet=True)

    add_textbox(slide, driver_card.left + Inches(0.18), driver_card.top + Inches(1.68), Inches(2.2), Inches(0.18), "OUTPUT TARGETS", size=10, color_hex=TEXT_SUB, bold=True, font_name=FONT_LATIN)
    add_textbox(slide, driver_card.left + Inches(0.18), driver_card.top + Inches(1.98), driver_card.width - Inches(0.36), Inches(0.20), "일별 시나리오 상태 | 주요 근거 지표", size=9, color_hex=TEXT_SUB, font_name=FONT_LATIN)
    add_footer(slide, 4)


def slide_roadmap(prs: Presentation, sections, assets) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, assets["roadmap_bg"])

    phases = extract_phases(sections["slide_5"]["lines"])

    add_textbox(slide, gx(0), gy(0.48), gw(5.8), Inches(0.30), "Phase 1-10 Roadmap", size=24, color_hex=TEXT, bold=True, font_name=FONT_LATIN)
    add_textbox(slide, gx(0), gy(0.88), gw(4.6), Inches(0.24), "Current focus: Explainable summary layer (Phase 4)", size=12, color_hex=GREEN, bold=True, font_name=FONT_LATIN)

    rail = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, gx(0), gy(1.30), gw(9.8), Inches(0.16))
    rail.fill.solid()
    rail.fill.fore_color.rgb = rgb("243149")
    rail.line.fill.background()
    progress = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, gx(0), gy(1.30), Inches(4.22), Inches(0.16))
    progress.fill.solid()
    progress.fill.fore_color.rgb = rgb(GREEN)
    progress.fill.transparency = 0.24
    progress.line.fill.background()
    add_textbox(slide, gx(9.96), gy(1.14), gw(1.8), Inches(0.24), "40% implemented", size=11, color_hex=TEXT, bold=True, font_name=FONT_LATIN)

    for idx, phase in enumerate(phases):
        row = idx // 5
        col = idx % 5
        x = gx(col * 2.42)
        y = gy(1.78 + row * 1.96)
        is_current = idx == 3
        accent = GREEN if is_current else MAGENTA if idx > 3 else TEXT_SUB
        card = add_glass_card(slide, x, y, Inches(2.34), Inches(1.54), fill_hex=CARD_ALT if is_current else CARD, transparency=0.12 if is_current else 0.20, line_hex=GREEN if is_current else CARD_STROKE)
        add_textbox(slide, card.left + Inches(0.16), card.top + Inches(0.16), Inches(0.58), Inches(0.20), f"{idx + 1:02d}", size=11, color_hex=accent, bold=True, font_name=FONT_LATIN)
        add_textbox(slide, card.left + Inches(0.16), card.top + Inches(0.44), card.width - Inches(0.32), Inches(0.42), str(phase["title"]).replace("Phase ", ""), size=13, color_hex=TEXT, bold=True)
        if idx == 6:
            body = "GFC·팬데믹·전쟁·금리 충격 검증"
        else:
            body = str(phase["bullets"][0]) if phase["bullets"] else ""
        add_textbox(slide, card.left + Inches(0.16), card.top + Inches(0.92), card.width - Inches(0.32), Inches(0.40), body, size=10, color_hex=TEXT_SUB)

    validation_card = add_glass_card(slide, gx(0), gy(5.78), gw(12.0), Inches(0.94), fill_hex=CARD_ALT)
    add_textbox(slide, validation_card.left + Inches(0.18), validation_card.top + Inches(0.16), Inches(1.8), Inches(0.18), "PHASE 7 CHECKS", size=10, color_hex=GREEN, bold=True, font_name=FONT_LATIN)
    add_textbox(slide, validation_card.left + Inches(0.18), validation_card.top + Inches(0.42), validation_card.width - Inches(0.36), Inches(0.24), "GFC type | Pandemic shock | War / Energy | Rate shock | 역사적 구간 상태 전환 검증", size=10, color_hex=TEXT_SUB, font_name=FONT_LATIN)
    add_footer(slide, 5)


def slide_effect(prs: Presentation, sections, assets) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, assets["default_bg"])
    bullets = extract_bullets(sections["slide_6"]["lines"])

    add_textbox(slide, gx(0), gy(0.48), gw(5.8), Inches(0.30), "리스크 관리와의 연결 방식", size=24, color_hex=TEXT, bold=True)
    add_textbox(slide, gx(0), gy(0.88), gw(5.4), Inches(0.24), "Risk visibility first, hedge action second.", size=12, color_hex=GREEN, bold=True, font_name=FONT_LATIN)

    flow = add_glass_card(slide, gx(0), gy(1.42), gw(6.1), Inches(2.42))
    add_textbox(slide, flow.left + Inches(0.18), flow.top + Inches(0.18), Inches(2.0), Inches(0.18), "DOWNSTREAM FLOW", size=10, color_hex=TEXT_SUB, bold=True, font_name=FONT_LATIN)
    nodes = [
        ("Market State", "Regime detection\nand drivers", GREEN),
        ("Risk Overlay", "Interpret current\nmarket context", MAGENTA),
        ("Hedge Action", "Prioritize hedge\nrecommendations", CYAN),
    ]
    for idx, (title, body, accent) in enumerate(nodes):
        x = flow.left + Inches(0.20 + idx * 1.95)
        card = add_glass_card(slide, x, flow.top + Inches(0.58), Inches(1.72), Inches(1.38), fill_hex=CARD_ALT)
        add_textbox(slide, card.left + Inches(0.14), card.top + Inches(0.14), card.width - Inches(0.28), Inches(0.22), title, size=12, color_hex=accent, bold=True, font_name=FONT_LATIN)
        add_textbox(slide, card.left + Inches(0.14), card.top + Inches(0.46), card.width - Inches(0.28), Inches(0.60), body, size=10, color_hex=TEXT)

    metric_y = gy(1.42)
    add_kpi_card(slide, gx(6.5), metric_y, Inches(1.72), Inches(1.16), "CVaR", "Tail-risk", "future overlay", GREEN)
    add_kpi_card(slide, gx(8.42), metric_y, Inches(1.72), Inches(1.16), "MDD", "Drawdown", "guardrail", MAGENTA)
    add_kpi_card(slide, gx(10.34), metric_y, Inches(1.72), Inches(1.16), "Sharpe", "Filter", "risk-adjusted lens", CYAN)

    hedge_card = add_glass_card(slide, gx(6.5), gy(2.74), gw(5.5), Inches(1.10), fill_hex=CARD_ALT)
    add_textbox(slide, hedge_card.left + Inches(0.18), hedge_card.top + Inches(0.18), Inches(2.8), Inches(0.18), "ILLUSTRATIVE HEDGE RECOMMENDATION PANEL", size=10, color_hex=TEXT_SUB, bold=True, font_name=FONT_LATIN)
    add_textbox(slide, hedge_card.left + Inches(0.18), hedge_card.top + Inches(0.48), hedge_card.width - Inches(0.36), Inches(0.34), "시장 상태 엔진이 완성되면, CVaR·MDD·Sharpe 같은 리스크 카드 위에 상황별 오버레이 우선순위를 자연스럽게 연결할 수 있습니다.", size=11, color_hex=TEXT)

    effect_card = add_glass_card(slide, gx(0), gy(4.10), gw(12.0), Inches(2.18))
    add_textbox(slide, effect_card.left + Inches(0.18), effect_card.top + Inches(0.18), Inches(2.0), Inches(0.18), "EXPECTED EFFECT", size=10, color_hex=GREEN, bold=True, font_name=FONT_LATIN)
    add_paragraph_block(slide, effect_card.left + Inches(0.18), effect_card.top + Inches(0.48), effect_card.width - Inches(0.36), Inches(1.24), bullets[:4], size=12, color_hex=TEXT, bullet=True)
    add_footer(slide, 6)


def slide_research(prs: Presentation, assets) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, assets["analysis_bg"])

    add_textbox(slide, gx(0), gy(0.48), gw(5.8), Inches(0.30), "Research Foundation", size=24, color_hex=TEXT, bold=True, font_name=FONT_LATIN)
    add_textbox(slide, gx(0), gy(0.88), gw(6.0), Inches(0.24), "Appendix references distilled into practical design signals.", size=12, color_hex=GREEN, bold=True, font_name=FONT_LATIN)

    cards = [
        ("AI in Quant Survey", "정형 + 비정형 + 자동화 구조로 시장 분석이 진화 중", GREEN),
        ("Mixture of Experts", "시장 상태가 다르면 단일 모델보다 적응형 구조가 유리", MAGENTA),
        ("Regime Switching", "상태 전환을 통계적으로 다루는 고전적 기반", CYAN),
        ("Open Source Stack", "Qlib, FinGPT, FinRL-X로 실무형 확장 가능", AMBER),
    ]
    for idx, (title, body, accent) in enumerate(cards):
        row = idx // 2
        col = idx % 2
        x = gx(col * 6.1)
        y = gy(1.42 + row * 2.20)
        card = add_glass_card(slide, x, y, gw(5.7), Inches(1.78), fill_hex=CARD_ALT if idx in (1, 2) else CARD)
        add_textbox(slide, card.left + Inches(0.18), card.top + Inches(0.18), card.width - Inches(0.36), Inches(0.24), title, size=16, color_hex=accent, bold=True, font_name=FONT_LATIN)
        add_textbox(slide, card.left + Inches(0.18), card.top + Inches(0.58), card.width - Inches(0.36), Inches(0.60), body, size=12, color_hex=TEXT)

    ref_card = add_glass_card(slide, gx(0), gy(5.98), gw(12.0), Inches(0.86))
    add_textbox(slide, ref_card.left + Inches(0.18), ref_card.top + Inches(0.16), Inches(2.0), Inches(0.18), "REFERENCE SET", size=10, color_hex=TEXT_SUB, bold=True, font_name=FONT_LATIN)
    refs = [
        "From Deep Learning to LLMs (2025)",
        "Adaptive Market Intelligence (2025)",
        "MIGA (2024)",
        "statsmodels MarkovRegression",
        "Qlib / FinGPT / FinRL-X",
    ]
    add_textbox(slide, ref_card.left + Inches(0.18), ref_card.top + Inches(0.42), ref_card.width - Inches(0.36), Inches(0.24), " | ".join(refs), size=10, color_hex=TEXT_SUB, font_name=FONT_LATIN)
    add_footer(slide, 7)


def build_deck(run_id: str | None = None) -> Path:
    ensure_dirs()
    sections = parse_markdown_sections(PLAN_PATH)
    data = load_data(run_id=run_id)

    assets = {
        "hero_bg": ASSET_DIR / "bg_hero.png",
        "analysis_bg": ASSET_DIR / "bg_analysis.png",
        "roadmap_bg": ASSET_DIR / "bg_roadmap.png",
        "default_bg": ASSET_DIR / "bg_default.png",
        "bar_chart": ASSET_DIR / "bar_chart.png",
        "donut_chart": ASSET_DIR / "donut_chart.png",
        "history_chart": ASSET_DIR / "history_chart.png",
    }

    create_background(assets["hero_bg"], "hero")
    create_background(assets["analysis_bg"], "analysis")
    create_background(assets["roadmap_bg"], "roadmap")
    create_background(assets["default_bg"], "default")
    create_bar_chart(data["latest_rows"], assets["bar_chart"])
    create_donut_chart(data["latest_rows"], assets["donut_chart"])
    create_history_chart(data["state_rows"], data["latest_rows"], assets["history_chart"])

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slide_cover(prs, sections, data, assets)
    slide_need(prs, sections, assets)
    slide_engine(prs, data, assets)
    slide_output(prs, sections, data, assets)
    slide_roadmap(prs, sections, assets)
    slide_effect(prs, sections, assets)
    slide_research(prs, assets)

    prs.save(PPT_PATH)
    return PPT_PATH


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate the market-state research dashboard deck.")
    parser.add_argument("--run-id", default=None, help="Use one specific pipeline run id instead of independently picking latest files.")
    args = parser.parse_args()
    output = build_deck(run_id=args.run_id)
    print(output)
